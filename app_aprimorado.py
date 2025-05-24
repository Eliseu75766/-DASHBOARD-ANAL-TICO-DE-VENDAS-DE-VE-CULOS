import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime, timedelta
import io
import base64
from utils.chart_factory import *
from utils.ai_insights import *
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import tempfile
import os
import re
from fpdf import FPDF
import json

# Configuração da página
st.set_page_config(
    page_title="Dashboard Analítico de Vendas de Veículos",
    page_icon="🚗",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Carregar CSS
with open('assets/css/style.css') as f:
    st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

# Função para criar link de download
def get_download_link(df, filename, link_text):
    csv = df.to_csv(index=False)
    b64 = base64.b64encode(csv.encode()).decode()
    href = f'<a href="data:file/csv;base64,{b64}" download="{filename}" class="download-link">{link_text}</a>'
    return href

# Função para exportar para PDF
def export_to_pdf(vendas, metas, insights_data, filtro_periodo=None, filtro_categorias=None, filtro_canais=None):
    pdf = FPDF()
    pdf.add_page()
    
    # Configurar fonte
    pdf.add_font('Montserrat', '', 'assets/fonts/Montserrat-Regular.ttf', uni=True)
    pdf.add_font('Montserrat', 'B', 'assets/fonts/Montserrat-Bold.ttf', uni=True)
    
    # Título
    pdf.set_font('Montserrat', 'B', 24)
    pdf.set_text_color(0, 255, 255)
    pdf.cell(0, 20, 'Dashboard Analítico de Vendas de Veículos', 0, 1, 'C')
    
    # Subtítulo
    pdf.set_font('Montserrat', '', 14)
    pdf.set_text_color(255, 95, 31)
    
    # Período
    if filtro_periodo:
        data_inicio, data_fim = filtro_periodo
        periodo_texto = f"Período: {data_inicio.strftime('%d/%m/%Y')} a {data_fim.strftime('%d/%m/%Y')}"
    else:
        periodo_texto = "Período: Todos os dados"
    
    pdf.cell(0, 10, periodo_texto, 0, 1, 'C')
    
    # Filtros aplicados
    pdf.set_font('Montserrat', '', 10)
    pdf.set_text_color(248, 248, 255)
    
    if filtro_categorias:
        categorias_texto = f"Categorias: {', '.join(filtro_categorias)}"
    else:
        categorias_texto = "Categorias: Todas"
    
    if filtro_canais:
        canais_texto = f"Canais: {', '.join(filtro_canais)}"
    else:
        canais_texto = "Canais: Todos"
    
    pdf.cell(0, 8, categorias_texto, 0, 1, 'C')
    pdf.cell(0, 8, canais_texto, 0, 1, 'C')
    
    pdf.ln(10)
    
    # Resumo dos dados
    pdf.set_font('Montserrat', 'B', 16)
    pdf.set_text_color(0, 255, 255)
    pdf.cell(0, 10, 'Resumo dos Dados', 0, 1, 'L')
    
    pdf.set_font('Montserrat', '', 12)
    pdf.set_text_color(248, 248, 255)
    
    # Criar tabela de resumo
    resumo = [
        ["Faturamento Total", f"R$ {insights_data['resumo_geral']['faturamento_total']:,.2f}"],
        ["Total de Vendas", f"{insights_data['resumo_geral']['total_vendas']:,}"],
        ["Lucro Total", f"R$ {insights_data['resumo_geral']['lucro_total']:,.2f}"],
        ["Margem Média", f"{insights_data['resumo_geral']['margem_media']:.2f}%"],
        ["Atingimento de Metas", f"{insights_data['metas']['atingimento_medio']:.1f}%"]
    ]
    
    col_width = 95
    row_height = 10
    
    for row in resumo:
        pdf.cell(col_width, row_height, row[0], 1, 0, 'L')
        pdf.cell(col_width, row_height, row[1], 1, 1, 'R')
    
    pdf.ln(10)
    
    # Insights principais
    pdf.set_font('Montserrat', 'B', 16)
    pdf.set_text_color(0, 255, 255)
    pdf.cell(0, 10, 'Insights Principais', 0, 1, 'L')
    
    pdf.set_font('Montserrat', '', 12)
    pdf.set_text_color(248, 248, 255)
    
    # Extrair texto limpo dos insights (remover tags HTML)
    narrativa = generate_narrative(insights_data)
    narrativa_limpa = re.sub('<.*?>', '', narrativa)
    
    # Quebrar texto em linhas
    pdf.multi_cell(0, 8, narrativa_limpa)
    
    pdf.ln(10)
    
    # Recomendações estratégicas
    pdf.set_font('Montserrat', 'B', 16)
    pdf.set_text_color(0, 255, 255)
    pdf.cell(0, 10, 'Recomendações Estratégicas', 0, 1, 'L')
    
    pdf.set_font('Montserrat', '', 12)
    pdf.set_text_color(248, 248, 255)
    
    # Extrair texto limpo das recomendações (remover tags HTML)
    recomendacoes = generate_strategic_recommendations(insights_data)
    recomendacoes_limpas = re.sub('<.*?>', '', recomendacoes)
    
    # Quebrar texto em linhas
    pdf.multi_cell(0, 8, recomendacoes_limpas)
    
    # Rodapé
    pdf.set_y(-15)
    pdf.set_font('Montserrat', 'I', 8)
    pdf.set_text_color(248, 248, 255)
    pdf.cell(0, 10, f'Gerado em {datetime.now().strftime("%d/%m/%Y %H:%M:%S")}', 0, 0, 'C')
    
    # Retornar PDF como buffer
    pdf_buffer = io.BytesIO()
    pdf.output(pdf_buffer)
    pdf_buffer.seek(0)
    
    return pdf_buffer

# Função para exportar para PowerPoint
def export_to_ppt(vendas, metas, insights_data, filtro_periodo=None, filtro_categorias=None, filtro_canais=None):
    # Criar apresentação
    prs = Presentation()
    
    # Slide de título
    slide_layout = prs.slide_layouts[0]  # Layout de título
    slide = prs.slides.add_slide(slide_layout)
    
    # Configurar título
    title = slide.shapes.title
    title.text = "Dashboard Analítico de Vendas de Veículos"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 255)
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    # Configurar subtítulo
    subtitle = slide.placeholders[1]
    
    # Período
    if filtro_periodo:
        data_inicio, data_fim = filtro_periodo
        periodo_texto = f"Período: {data_inicio.strftime('%d/%m/%Y')} a {data_fim.strftime('%d/%m/%Y')}"
    else:
        periodo_texto = "Período: Todos os dados"
    
    # Filtros aplicados
    if filtro_categorias:
        categorias_texto = f"Categorias: {', '.join(filtro_categorias)}"
    else:
        categorias_texto = "Categorias: Todas"
    
    if filtro_canais:
        canais_texto = f"Canais: {', '.join(filtro_canais)}"
    else:
        canais_texto = "Canais: Todos"
    
    subtitle.text = f"{periodo_texto}\n{categorias_texto}\n{canais_texto}"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 95, 31)
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Slide de resumo
    slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
    slide = prs.slides.add_slide(slide_layout)
    
    # Configurar título
    title = slide.shapes.title
    title.text = "Resumo dos Dados"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 255)
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    # Criar tabela de resumo
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)
    shape = slide.shapes.add_table(6, 2, x, y, cx, cy)
    table = shape.table
    
    # Cabeçalho
    table.cell(0, 0).text = "Métrica"
    table.cell(0, 1).text = "Valor"
    
    # Dados
    table.cell(1, 0).text = "Faturamento Total"
    table.cell(1, 1).text = f"R$ {insights_data['resumo_geral']['faturamento_total']:,.2f}"
    
    table.cell(2, 0).text = "Total de Vendas"
    table.cell(2, 1).text = f"{insights_data['resumo_geral']['total_vendas']:,}"
    
    table.cell(3, 0).text = "Lucro Total"
    table.cell(3, 1).text = f"R$ {insights_data['resumo_geral']['lucro_total']:,.2f}"
    
    table.cell(4, 0).text = "Margem Média"
    table.cell(4, 1).text = f"{insights_data['resumo_geral']['margem_media']:.2f}%"
    
    table.cell(5, 0).text = "Atingimento de Metas"
    table.cell(5, 1).text = f"{insights_data['metas']['atingimento_medio']:.1f}%"
    
    # Slide de insights
    slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
    slide = prs.slides.add_slide(slide_layout)
    
    # Configurar título
    title = slide.shapes.title
    title.text = "Insights Principais"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 255)
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    # Extrair texto limpo dos insights (remover tags HTML)
    narrativa = generate_narrative(insights_data)
    narrativa_limpa = re.sub('<.*?>', '', narrativa)
    
    # Adicionar texto
    content = slide.placeholders[1]
    content.text = narrativa_limpa
    
    # Slide de recomendações
    slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
    slide = prs.slides.add_slide(slide_layout)
    
    # Configurar título
    title = slide.shapes.title
    title.text = "Recomendações Estratégicas"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 255)
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    # Extrair texto limpo das recomendações (remover tags HTML)
    recomendacoes = generate_strategic_recommendations(insights_data)
    recomendacoes_limpas = re.sub('<.*?>', '', recomendacoes)
    
    # Adicionar texto
    content = slide.placeholders[1]
    content.text = recomendacoes_limpas
    
    # Retornar PPT como buffer
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    return ppt_buffer

# Função para carregar dados
@st.cache_data
def load_data():
    # Carregar dados de vendas
    vendas = pd.read_csv('data/vendas.csv')
    vendas['data_venda'] = pd.to_datetime(vendas['data_venda'])
    vendas['periodo'] = vendas['data_venda'].dt.strftime('%Y-%m')
    
    # Extrair hora da data_venda para análise por horário
    vendas['hora'] = vendas['data_venda'].dt.hour
    
    # Carregar dados de metas
    metas = pd.read_csv('data/metas.csv')
    
    # Carregar dados de modelos
    modelos = pd.read_csv('data/modelos.csv')
    
    return vendas, metas, modelos

# Carregar dados
vendas, metas, modelos = load_data()

# Verificar se os dados foram carregados corretamente
if vendas.empty or metas.empty or modelos.empty:
    st.error("Erro ao carregar os dados. Verifique os arquivos CSV.")
    st.stop()

# Sidebar para filtros
with st.sidebar:
    st.markdown('<div class="sidebar-title">Filtros de Análise</div>', unsafe_allow_html=True)
    
    st.markdown('<hr>', unsafe_allow_html=True)
    
    # Filtro de período
    st.markdown('<div class="filter-title">Período de Análise</div>', unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        data_min = vendas['data_venda'].min().date()
        data_max = vendas['data_venda'].max().date()
        
        data_inicio = st.date_input(
            "De",
            value=data_min,
            min_value=data_min,
            max_value=data_max,
            format="YYYY/MM/DD"
        )
    
    with col2:
        data_fim = st.date_input(
            "Até",
            value=data_max,
            min_value=data_min,
            max_value=data_max,
            format="YYYY/MM/DD"
        )
    
    # Converter para datetime
    data_inicio = pd.to_datetime(data_inicio)
    data_fim = pd.to_datetime(data_fim)
    
    # Filtro de categorias
    st.markdown('<div class="filter-title">Categorias de Veículos</div>', unsafe_allow_html=True)
    
    categorias = sorted(vendas['categoria'].unique())
    
    categorias_selecionadas = st.multiselect(
        "Selecione as categorias",
        options=categorias,
        default=categorias[:3] if len(categorias) > 3 else categorias
    )
    
    # Filtro de canais de venda
    st.markdown('<div class="filter-title">Canais de Venda</div>', unsafe_allow_html=True)
    
    canais = sorted(vendas['canal_venda'].unique())
    
    canais_selecionados = st.multiselect(
        "Selecione os canais",
        options=canais,
        default=canais[:3] if len(canais) > 3 else canais
    )
    
    # Botão para aplicar filtros
    filtros_aplicados = st.button("Aplicar Filtros", type="primary")
    
    # Informações sobre o dashboard
    st.markdown('<div class="sidebar-info">Sobre o Dashboard</div>', unsafe_allow_html=True)
    st.markdown("""
    Dashboard analítico interativo para visualização de dados de vendas de veículos, com foco em data storytelling e insights acionáveis.
    
    Desenvolvido com Streamlit, Plotly e IA para análise avançada de dados.
    """)

# Aplicar filtros
if filtros_aplicados or 'filtro_aplicado' not in st.session_state:
    st.session_state.filtro_aplicado = True
    
    # Filtro de período
    filtro_periodo = (data_inicio, data_fim)
    
    # Filtro de categorias
    filtro_categorias = categorias_selecionadas if categorias_selecionadas else None
    
    # Filtro de canais
    filtro_canais = canais_selecionados if canais_selecionados else None
else:
    # Usar filtros da sessão anterior
    filtro_periodo = (st.session_state.get('data_inicio', data_inicio), st.session_state.get('data_fim', data_fim))
    filtro_categorias = st.session_state.get('categorias_selecionadas', categorias_selecionadas if categorias_selecionadas else None)
    filtro_canais = st.session_state.get('canais_selecionados', canais_selecionados if canais_selecionados else None)

# Salvar filtros na sessão
st.session_state.data_inicio = filtro_periodo[0]
st.session_state.data_fim = filtro_periodo[1]
st.session_state.categorias_selecionadas = filtro_categorias
st.session_state.canais_selecionados = filtro_canais

# Título principal
st.markdown(
    """
    <div class="main-title-container">
        <div class="main-title">
            <span class="icon">🚗</span> DASHBOARD ANALÍTICO DE<br>VENDAS DE VEÍCULOS
        </div>
        <div class="subtitle">Análise de Performance 2024-2025</div>
    </div>
    """, 
    unsafe_allow_html=True
)

# Tabs para navegação
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "📊 Faturamento", 
    "💰 Margem de Lucro", 
    "💵 Ticket Médio", 
    "📅 Análise Mensal", 
    "🕒 Análise por Horário", 
    "🧠 IA Insights"
])

# Gerar insights avançados
insights_data = generate_advanced_insights(vendas, metas, modelos, filtro_periodo, filtro_categorias, filtro_canais)

# Tab 1: Faturamento
with tab1:
    st.markdown('<div class="tab-title">Análise de Faturamento</div>', unsafe_allow_html=True)
    
    # KPIs
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        # Faturamento total
        faturamento_total = insights_data['resumo_geral']['faturamento_total']
        
        st.markdown('<div class="kpi-card">', unsafe_allow_html=True)
        st.markdown('<div class="kpi-title">Faturamento Total</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="kpi-value">R$ {faturamento_total:,.2f}</div>', unsafe_allow_html=True)
        
        # Tendência
        if insights_data['resumo_geral']['tendencia_crescimento'] is not None:
            tendencia = insights_data['resumo_geral']['tendencia_crescimento']
            icon = "↑" if tendencia > 0 else "↓"
            color = "positive" if tendencia > 0 else "negative"
            st.markdown(f'<div class="kpi-trend {color}">{icon} {abs(tendencia):.1f}%</div>', unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col2:
        # Quantidade de vendas
        qtd_vendas = insights_data['resumo_geral']['total_vendas']
        
        st.markdown('<div class="kpi-card">', unsafe_allow_html=True)
        st.markdown('<div class="kpi-title">Quantidade de Vendas</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="kpi-value">{qtd_vendas:,}</div>', unsafe_allow_html=True)
        
        # Tendência
        if insights_data['tendencias']['volume'] is not None:
            tendencia = insights_data['tendencias']['volume']
            icon = "↑" if tendencia > 0 else "↓"
            color = "positive" if tendencia > 0 else "negative"
            st.markdown(f'<div class="kpi-trend {color}">{icon} {abs(tendencia):.1f}%</div>', unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col3:
        # Ticket médio
        if qtd_vendas > 0:
            ticket_medio = faturamento_total / qtd_vendas
        else:
            ticket_medio = 0
        
        st.markdown('<div class="kpi-card">', unsafe_allow_html=True)
        st.markdown('<div class="kpi-title">Ticket Médio</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="kpi-value">R$ {ticket_medio:,.2f}</div>', unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col4:
        # Margem média
        margem_media = insights_data['resumo_geral']['margem_media']
        
        st.markdown('<div class="kpi-card">', unsafe_allow_html=True)
        st.markdown('<div class="kpi-title">Margem Média</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="kpi-value">{margem_media:.2f}%</div>', unsafe_allow_html=True)
        
        # Tendência
        if insights_data['tendencias']['margem'] is not None:
            tendencia = insights_data['tendencias']['margem']
            icon = "↑" if tendencia > 0 else "↓"
            color = "positive" if tendencia > 0 else "negative"
            st.markdown(f'<div class="kpi-trend {color}">{icon} {abs(tendencia):.1f}%</div>', unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    # Gráfico de faturamento
    st.markdown('<div class="chart-container">', unsafe_allow_html=True)
    
    fig_faturamento, faturamento_mensal = create_faturamento_chart(vendas, metas, filtro_periodo, filtro_categorias, filtro_canais)
    st.plotly_chart(fig_faturamento, use_container_width=True)
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Análise por categoria
    st.markdown('<div class="section-title">Análise por Categoria</div>', unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        # Agrupar por categoria
        categoria_stats = vendas.copy()
        if filtro_periodo:
            data_inicio, data_fim = filtro_periodo
            categoria_stats = categoria_stats[(categoria_stats['data_venda'] >= data_inicio) & (categoria_stats['data_venda'] <= data_fim)]
        
        if filtro_categorias:
            categoria_stats = categoria_stats[categoria_stats['categoria'].isin(filtro_categorias)]
        
        if filtro_canais:
            categoria_stats = categoria_stats[categoria_stats['canal_venda'].isin(filtro_canais)]
        
        categoria_stats = categoria_stats.groupby('categoria').agg({
            'id_venda': 'count',
            'preco_venda': 'sum',
            'lucro': 'sum'
        }).reset_index()
        
        categoria_stats['margem'] = (categoria_stats['lucro'] / categoria_stats['preco_venda']) * 100
        categoria_stats['ticket_medio'] = categoria_stats['preco_venda'] / categoria_stats['id_venda']
        
        # Ordenar por faturamento
        categoria_stats = categoria_stats.sort_values('preco_venda', ascending=False)
        
        # Criar gráfico de barras
        fig = go.Figure()
        
        # Cores para cada categoria
        cores = {
            'SUV': '#00FFFF',
            'Sedan': '#FF5F1F',
            'Hatch': '#F8F8FF',
            'Pickup': '#9933FF',
            'Elétrico': '#33FF99'
        }
        
        # Adicionar barras para cada categoria
        fig.add_trace(
            go.Bar(
                x=categoria_stats['categoria'],
                y=categoria_stats['preco_venda'],
                name="Faturamento",
                marker=dict(color=[cores.get(cat, '#F8F8FF') for cat in categoria_stats['categoria']]),
                hovertemplate='<b>%{x}</b><br>Faturamento: R$ %{y:,.2f}<br>Quantidade: %{customdata[0]:,}<br>Margem: %{customdata[1]:.2f}%<extra></extra>',
                customdata=np.column_stack((categoria_stats['id_venda'], categoria_stats['margem'])),
                text=categoria_stats['preco_venda'].apply(lambda x: f'R$ {x:,.0f}'),
                textposition='inside'
            )
        )
        
        # Personalizar layout
        fig.update_layout(
            title={
                'text': "Faturamento por Categoria",
                'y':0.95,
                'x':0.5,
                'xanchor': 'center',
                'yanchor': 'top',
                'font': dict(family="Orbitron", size=20, color="#F8F8FF")
            },
            paper_bgcolor='rgba(13, 13, 13, 0.0)',
            plot_bgcolor='rgba(13, 13, 13, 0.0)',
            margin=dict(l=20, r=20, t=80, b=20),
            height=400,
            xaxis=dict(
                title=dict(
                    text="Categoria",
                    font=dict(family="Montserrat", color="#F8F8FF")
                ),
                tickfont=dict(family="Montserrat", color="#F8F8FF"),
                showgrid=False,
                zeroline=False
            ),
            yaxis=dict(
                title=dict(
                    text="Faturamento (R$)",
                    font=dict(family="Montserrat", color="#F8F8FF")
                ),
                tickfont=dict(family="Montserrat", color="#F8F8FF"),
                showgrid=True,
                gridcolor='rgba(248, 248, 255, 0.1)',
                zeroline=False
            )
        )
        
        st.plotly_chart(fig, use_container_width=True)
    
    with col2:
        # Criar gráfico de pizza para distribuição de vendas por categoria
        fig = go.Figure()
        
        # Adicionar pizza
        fig.add_trace(
            go.Pie(
                labels=categoria_stats['categoria'],
                values=categoria_stats['id_venda'],
                hole=0.5,
                marker=dict(colors=[cores.get(cat, '#F8F8FF') for cat in categoria_stats['categoria']]),
                textinfo='label+percent',
                hovertemplate='<b>%{label}</b><br>Quantidade: %{value:,}<br>Percentual: %{percent}<extra></extra>',
                textfont=dict(family="Montserrat", color="#0D0D0D", size=12),
                insidetextorientation='radial'
            )
        )
        
        # Personalizar layout
        fig.update_layout(
            title={
                'text': "Distribuição de Vendas por Categoria",
                'y':0.95,
                'x':0.5,
                'xanchor': 'center',
                'yanchor': 'top',
                'font': dict(family="Orbitron", size=20, color="#F8F8FF")
            },
            paper_bgcolor='rgba(13, 13, 13, 0.0)',
            plot_bgcolor='rgba(13, 13, 13, 0.0)',
            margin=dict(l=20, r=20, t=80, b=20),
            height=400,
            legend=dict(
                orientation="h",
                yanchor="bottom",
                y=-0.2,
                xanchor="center",
                x=0.5,
                font=dict(family="Montserrat", color="#F8F8FF")
            )
        )
        
        st.plotly_chart(fig, use_container_width=True)
    
    # Análise por canal
    st.markdown('<div class="section-title">Análise por Canal de Venda</div>', unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        # Agrupar por canal
        canal_stats = vendas.copy()
        if filtro_periodo:
            data_inicio, data_fim = filtro_periodo
            canal_stats = canal_stats[(canal_stats['data_venda'] >= data_inicio) & (canal_stats['data_venda'] <= data_fim)]
        
        if filtro_categorias:
            canal_stats = canal_stats[canal_stats['categoria'].isin(filtro_categorias)]
        
        if filtro_canais:
            canal_stats = canal_stats[canal_stats['canal_venda'].isin(filtro_canais)]
        
        canal_stats = canal_stats.groupby('canal_venda').agg({
            'id_venda': 'count',
            'preco_venda': 'sum',
            'lucro': 'sum'
        }).reset_index()
        
        canal_stats['margem'] = (canal_stats['lucro'] / canal_stats['preco_venda']) * 100
        canal_stats['ticket_medio'] = canal_stats['preco_venda'] / canal_stats['id_venda']
        
        # Ordenar por faturamento
        canal_stats = canal_stats.sort_values('preco_venda', ascending=False)
        
        # Criar gráfico de barras
        fig = go.Figure()
        
        # Cores para cada canal
        cores_canal = {
            'Showroom': '#00FFFF',
            'Online': '#FF5F1F',
            'Concessionária': '#F8F8FF',
            'Parceiro': '#9933FF'
        }
        
        # Adicionar barras para cada canal
        fig.add_trace(
            go.Bar(
                x=canal_stats['canal_venda'],
                y=canal_stats['preco_venda'],
                name="Faturamento",
                marker=dict(color=[cores_canal.get(canal, '#F8F8FF') for canal in canal_stats['canal_venda']]),
                hovertemplate='<b>%{x}</b><br>Faturamento: R$ %{y:,.2f}<br>Quantidade: %{customdata[0]:,}<br>Margem: %{customdata[1]:.2f}%<extra></extra>',
                customdata=np.column_stack((canal_stats['id_venda'], canal_stats['margem'])),
                text=canal_stats['preco_venda'].apply(lambda x: f'R$ {x:,.0f}'),
                textposition='inside'
            )
        )
        
        # Personalizar layout
        fig.update_layout(
            title={
                'text': "Faturamento por Canal de Venda",
                'y':0.95,
                'x':0.5,
                'xanchor': 'center',
                'yanchor': 'top',
                'font': dict(family="Orbitron", size=20, color="#F8F8FF")
            },
            paper_bgcolor='rgba(13, 13, 13, 0.0)',
            plot_bgcolor='rgba(13, 13, 13, 0.0)',
            margin=dict(l=20, r=20, t=80, b=20),
            height=400,
            xaxis=dict(
                title=dict(
                    text="Canal de Venda",
                    font=dict(family="Montserrat", color="#F8F8FF")
                ),
                tickfont=dict(family="Montserrat", color="#F8F8FF"),
                showgrid=False,
                zeroline=False
            ),
            yaxis=dict(
                title=dict(
                    text="Faturamento (R$)",
                    font=dict(family="Montserrat", color="#F8F8FF")
                ),
                tickfont=dict(family="Montserrat", color="#F8F8FF"),
                showgrid=True,
                gridcolor='rgba(248, 248, 255, 0.1)',
                zeroline=False
            )
        )
        
        st.plotly_chart(fig, use_container_width=True)
    
    with col2:
        # Criar gráfico de pizza para distribuição de vendas por canal
        fig = go.Figure()
        
        # Adicionar pizza
        fig.add_trace(
            go.Pie(
                labels=canal_stats['canal_venda'],
                values=canal_stats['id_venda'],
                hole=0.5,
                marker=dict(colors=[cores_canal.get(canal, '#F8F8FF') for canal in canal_stats['canal_venda']]),
                textinfo='label+percent',
                hovertemplate='<b>%{label}</b><br>Quantidade: %{value:,}<br>Percentual: %{percent}<extra></extra>',
                textfont=dict(family="Montserrat", color="#0D0D0D", size=12),
                insidetextorientation='radial'
            )
        )
        
        # Personalizar layout
        fig.update_layout(
            title={
                'text': "Distribuição de Vendas por Canal",
                'y':0.95,
                'x':0.5,
                'xanchor': 'center',
                'yanchor': 'top',
                'font': dict(family="Orbitron", size=20, color="#F8F8FF")
            },
            paper_bgcolor='rgba(13, 13, 13, 0.0)',
            plot_bgcolor='rgba(13, 13, 13, 0.0)',
            margin=dict(l=20, r=20, t=80, b=20),
            height=400,
            legend=dict(
                orientation="h",
                yanchor="bottom",
                y=-0.2,
                xanchor="center",
                x=0.5,
                font=dict(family="Montserrat", color="#F8F8FF")
            )
        )
        
        st.plotly_chart(fig, use_container_width=True)
    
    # Download de dados
    st.markdown('<div class="download-section">', unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown(get_download_link(faturamento_mensal, "faturamento_mensal.csv", "📥 Baixar Dados de Faturamento"), unsafe_allow_html=True)
    
    with col2:
        st.markdown(get_download_link(categoria_stats, "categoria_stats.csv", "📥 Baixar Dados por Categoria"), unsafe_allow_html=True)
    
    with col3:
        st.markdown(get_download_link(canal_stats, "canal_stats.csv", "📥 Baixar Dados por Canal"), unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)

# Tab 2: Margem de Lucro
with tab2:
    st.markdown('<div class="tab-title">Análise de Margem de Lucro</div>', unsafe_allow_html=True)
    
    # KPIs
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        # Lucro total
        lucro_total = insights_data['resumo_geral']['lucro_total']
        
        st.markdown('<div class="kpi-card">', unsafe_allow_html=True)
        st.markdown('<div class="kpi-title">Lucro Total</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="kpi-value">R$ {lucro_total:,.2f}</div>', unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col2:
        # Margem média
        margem_media = insights_data['resumo_geral']['margem_media']
        
        st.markdown('<div class="kpi-card">', unsafe_allow_html=True)
        st.markdown('<div class="kpi-title">Margem Média</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="kpi-value">{margem_media:.2f}%</div>', unsafe_allow_html=True)
        
        # Tendência
        if insights_data['tendencias']['margem'] is not None:
            tendencia = insights_data['tendencias']['margem']
            icon = "↑" if tendencia > 0 else "↓"
            color = "positive" if tendencia > 0 else "negative"
            st.markdown(f'<div class="kpi-trend {color}">{icon} {abs(tendencia):.1f}%</div>', unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col3:
        # Categoria mais lucrativa
        if insights_data['categorias']['mais_lucrativa'] is not None:
            categoria = insights_data['categorias']['mais_lucrativa']['categoria']
            margem = insights_data['categorias']['mais_lucrativa']['margem']
            
            st.markdown('<div class="kpi-card">', unsafe_allow_html=True)
            st.markdown('<div class="kpi-title">Categoria Mais Lucrativa</div>', unsafe_allow_html=True)
            st.markdown(f'<div class="kpi-value">{categoria}</div>', unsafe_allow_html=True)
            st.markdown(f'<div class="kpi-subtitle">Margem: {margem:.2f}%</div>', unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)
    
    with col4:
        # Canal mais lucrativo
        if insights_data['canais']['mais_lucrativo'] is not None:
            canal = insights_data['canais']['mais_lucrativo']['canal_venda']
            margem = insights_data['canais']['mais_lucrativo']['margem']
            
            st.markdown('<div class="kpi-card">', unsafe_allow_html=True)
            st.markdown('<div class="kpi-title">Canal Mais Lucrativo</div>', unsafe_allow_html=True)
            st.markdown(f'<div class="kpi-value">{canal}</div>', unsafe_allow_html=True)
            st.markdown(f'<div class="kpi-subtitle">Margem: {margem:.2f}%</div>', unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)
    
    # Gráfico de margem por categoria
    st.markdown('<div class="chart-container">', unsafe_allow_html=True)
    
    fig_margem, margem_categoria = create_margem_chart(vendas, filtro_periodo, filtro_categorias, filtro_canais)
    st.plotly_chart(fig_margem, use_container_width=True)
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Margem por canal
    st.markdown('<div class="section-title">Margem de Lucro por Canal de Vendas</div>', unsafe_allow_html=True)
    
    fig_margem_canal, margem_canal = create_margem_canal_chart(vendas, filtro_periodo, filtro_categorias, filtro_canais)
    st.plotly_chart(fig_margem_canal, use_container_width=True)
    
    # Análise detalhada por modelo
    st.markdown('<div class="section-title">Análise Detalhada por Modelo</div>', unsafe_allow_html=True)
    
    # Agrupar por modelo
    modelo_stats = vendas.copy()
    if filtro_periodo:
        data_inicio, data_fim = filtro_periodo
        modelo_stats = modelo_stats[(modelo_stats['data_venda'] >= data_inicio) & (modelo_stats['data_venda'] <= data_fim)]
    
    if filtro_categorias:
        modelo_stats = modelo_stats[modelo_stats['categoria'].isin(filtro_categorias)]
    
    if filtro_canais:
        modelo_stats = modelo_stats[modelo_stats['canal_venda'].isin(filtro_canais)]
    
    modelo_stats = modelo_stats.groupby('modelo').agg({
        'id_venda': 'count',
        'preco_venda': 'sum',
        'custo': 'sum',
        'lucro': 'sum'
    }).reset_index()
    
    modelo_stats['margem'] = (modelo_stats['lucro'] / modelo_stats['preco_venda']) * 100
    modelo_stats['ticket_medio'] = modelo_stats['preco_venda'] / modelo_stats['id_venda']
    
    # Ordenar por margem
    modelo_stats = modelo_stats.sort_values('margem', ascending=False)
    
    # Criar gráfico de barras
    fig = go.Figure()
    
    # Adicionar barras para cada modelo
    fig.add_trace(
        go.Bar(
            x=modelo_stats['modelo'],
            y=modelo_stats['margem'],
            name="Margem (%)",
            marker=dict(
                color=modelo_stats['margem'],
                colorscale=[
                    [0, 'rgba(255, 95, 31, 0.7)'],
                    [0.5, 'rgba(248, 248, 255, 0.7)'],
                    [1, 'rgba(0, 255, 255, 0.7)']
                ]
            ),
            hovertemplate='<b>%{x}</b><br>Margem: %{y:.2f}%<br>Lucro: R$ %{customdata[0]:,.2f}<br>Vendas: %{customdata[1]:,}<extra></extra>',
            customdata=np.column_stack((modelo_stats['lucro'], modelo_stats['id_venda'])),
            text=modelo_stats['margem'].apply(lambda x: f'{x:.2f}%'),
            textposition='inside'
        )
    )
    
    # Personalizar layout
    fig.update_layout(
        title={
            'text': "Margem de Lucro por Modelo",
            'y':0.95,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': dict(family="Orbitron", size=24, color="#F8F8FF")
        },
        paper_bgcolor='rgba(13, 13, 13, 0.0)',
        plot_bgcolor='rgba(13, 13, 13, 0.0)',
        margin=dict(l=20, r=20, t=80, b=20),
        height=500,
        xaxis=dict(
            title=dict(
                text="Modelo",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            showgrid=False,
            zeroline=False
        ),
        yaxis=dict(
            title=dict(
                text="Margem (%)",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            showgrid=True,
            gridcolor='rgba(248, 248, 255, 0.1)',
            zeroline=False
        )
    )
    
    st.plotly_chart(fig, use_container_width=True)
    
    # Download de dados
    st.markdown('<div class="download-section">', unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown(get_download_link(margem_categoria, "margem_categoria.csv", "📥 Baixar Dados de Margem por Categoria"), unsafe_allow_html=True)
    
    with col2:
        st.markdown(get_download_link(margem_canal, "margem_canal.csv", "📥 Baixar Dados de Margem por Canal"), unsafe_allow_html=True)
    
    with col3:
        st.markdown(get_download_link(modelo_stats, "modelo_stats.csv", "📥 Baixar Dados por Modelo"), unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)

# Tab 3: Ticket Médio
with tab3:
    st.markdown('<div class="tab-title">Análise de Ticket Médio</div>', unsafe_allow_html=True)
    
    # Gráficos de ticket médio
    fig_gauge, fig_line, ticket_medio = create_ticket_chart(vendas, metas, filtro_periodo, filtro_categorias, filtro_canais)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.plotly_chart(fig_gauge, use_container_width=True)
    
    with col2:
        # Insights de ticket médio
        st.markdown('<div class="insight-box">', unsafe_allow_html=True)
        st.markdown('<div class="insight-title">Insights de Ticket Médio</div>', unsafe_allow_html=True)
        
        ticket_insights = generate_ticket_insights(vendas, filtro_periodo, filtro_categorias, filtro_canais)
        st.markdown(ticket_insights, unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    # Evolução do ticket médio
    st.plotly_chart(fig_line, use_container_width=True)
    
    # Ticket médio por categoria
    st.markdown('<div class="section-title">Ticket Médio por Categoria</div>', unsafe_allow_html=True)
    
    fig_ticket_categoria, ticket_categoria = create_ticket_categoria_chart(vendas, filtro_periodo, filtro_categorias, filtro_canais)
    st.plotly_chart(fig_ticket_categoria, use_container_width=True)
    
    # Ticket médio por canal
    st.markdown('<div class="section-title">Ticket Médio por Canal de Venda</div>', unsafe_allow_html=True)
    
    # Agrupar por canal
    canal_ticket = vendas.copy()
    if filtro_periodo:
        data_inicio, data_fim = filtro_periodo
        canal_ticket = canal_ticket[(canal_ticket['data_venda'] >= data_inicio) & (canal_ticket['data_venda'] <= data_fim)]
    
    if filtro_categorias:
        canal_ticket = canal_ticket[canal_ticket['categoria'].isin(filtro_categorias)]
    
    if filtro_canais:
        canal_ticket = canal_ticket[canal_ticket['canal_venda'].isin(filtro_canais)]
    
    canal_ticket = canal_ticket.groupby('canal_venda').agg({
        'id_venda': 'count',
        'preco_venda': 'sum'
    }).reset_index()
    
    canal_ticket['ticket_medio'] = canal_ticket['preco_venda'] / canal_ticket['id_venda']
    
    # Ordenar por ticket médio
    canal_ticket = canal_ticket.sort_values('ticket_medio', ascending=False)
    
    # Criar gráfico de barras
    fig = go.Figure()
    
    # Cores para cada canal
    cores_canal = {
        'Showroom': '#00FFFF',
        'Online': '#FF5F1F',
        'Concessionária': '#F8F8FF',
        'Parceiro': '#9933FF'
    }
    
    # Adicionar barras para cada canal
    fig.add_trace(
        go.Bar(
            x=canal_ticket['canal_venda'],
            y=canal_ticket['ticket_medio'],
            name="Ticket Médio",
            marker=dict(color=[cores_canal.get(canal, '#F8F8FF') for canal in canal_ticket['canal_venda']]),
            hovertemplate='<b>%{x}</b><br>Ticket Médio: R$ %{y:,.2f}<br>Quantidade: %{customdata:,}<extra></extra>',
            customdata=canal_ticket['id_venda'],
            text=canal_ticket['ticket_medio'].apply(lambda x: f'R$ {x:,.2f}'),
            textposition='inside'
        )
    )
    
    # Personalizar layout
    fig.update_layout(
        title={
            'text': "Ticket Médio por Canal de Venda",
            'y':0.95,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': dict(family="Orbitron", size=24, color="#F8F8FF")
        },
        paper_bgcolor='rgba(13, 13, 13, 0.0)',
        plot_bgcolor='rgba(13, 13, 13, 0.0)',
        margin=dict(l=20, r=20, t=80, b=20),
        height=400,
        xaxis=dict(
            title=dict(
                text="Canal de Venda",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            showgrid=False,
            zeroline=False
        ),
        yaxis=dict(
            title=dict(
                text="Ticket Médio (R$)",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            showgrid=True,
            gridcolor='rgba(248, 248, 255, 0.1)',
            zeroline=False
        )
    )
    
    st.plotly_chart(fig, use_container_width=True)
    
    # Download de dados
    st.markdown('<div class="download-section">', unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown(get_download_link(ticket_medio, "ticket_medio.csv", "📥 Baixar Dados de Ticket Médio"), unsafe_allow_html=True)
    
    with col2:
        st.markdown(get_download_link(ticket_categoria, "ticket_categoria.csv", "📥 Baixar Dados de Ticket por Categoria"), unsafe_allow_html=True)
    
    with col3:
        st.markdown(get_download_link(canal_ticket, "canal_ticket.csv", "📥 Baixar Dados de Ticket por Canal"), unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)

# Tab 4: Análise Mensal
with tab4:
    st.markdown('<div class="tab-title">Análise Mensal</div>', unsafe_allow_html=True)
    
    # Heatmap de vendas por modelo e período
    st.markdown('<div class="chart-container">', unsafe_allow_html=True)
    
    fig_heatmap, matriz_vendas = create_heatmap(vendas, filtro_periodo, filtro_categorias, filtro_canais)
    st.plotly_chart(fig_heatmap, use_container_width=True)
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Análise de atingimento de metas
    st.markdown('<div class="section-title">Atingimento de Metas</div>', unsafe_allow_html=True)
    
    # Agrupar por período
    periodo_stats = vendas.copy()
    if filtro_periodo:
        data_inicio, data_fim = filtro_periodo
        periodo_stats = periodo_stats[(periodo_stats['data_venda'] >= data_inicio) & (periodo_stats['data_venda'] <= data_fim)]
    
    if filtro_categorias:
        periodo_stats = periodo_stats[periodo_stats['categoria'].isin(filtro_categorias)]
    
    if filtro_canais:
        periodo_stats = periodo_stats[periodo_stats['canal_venda'].isin(filtro_canais)]
    
    periodo_stats = periodo_stats.groupby('periodo').agg({
        'id_venda': 'count',
        'preco_venda': 'sum'
    }).reset_index()
    
    # Mesclar com metas
    metas_periodo = metas.copy()
    metas_periodo = pd.merge(periodo_stats, metas_periodo[['periodo', 'meta_faturamento']], on='periodo', how='left')
    
    # Calcular atingimento
    metas_periodo['atingimento'] = (metas_periodo['preco_venda'] / metas_periodo['meta_faturamento']) * 100
    
    # Criar gráfico de barras e linha
    fig = go.Figure()
    
    # Adicionar barras para faturamento
    fig.add_trace(
        go.Bar(
            x=metas_periodo['periodo'],
            y=metas_periodo['preco_venda'],
            name="Faturamento",
            marker=dict(color="#00FFFF"),
            hovertemplate='<b>%{x}</b><br>Faturamento: R$ %{y:,.2f}<br>Atingimento: %{customdata:.1f}%<extra></extra>',
            customdata=metas_periodo['atingimento'],
            text=metas_periodo['preco_venda'].apply(lambda x: f'R$ {x:,.0f}'),
            textposition='inside'
        )
    )
    
    # Adicionar linha para meta
    fig.add_trace(
        go.Scatter(
            x=metas_periodo['periodo'],
            y=metas_periodo['meta_faturamento'],
            name="Meta",
            line=dict(color="#FF5F1F", width=3, dash='dash'),
            mode='lines+markers',
            marker=dict(size=8, symbol='diamond', color="#FF5F1F"),
            hovertemplate='<b>%{x}</b><br>Meta: R$ %{y:,.2f}<extra></extra>'
        )
    )
    
    # Adicionar linha para atingimento
    fig.add_trace(
        go.Scatter(
            x=metas_periodo['periodo'],
            y=metas_periodo['atingimento'],
            name="Atingimento (%)",
            line=dict(color="#F8F8FF", width=3),
            mode='lines+markers+text',
            marker=dict(size=8, symbol='circle', color="#F8F8FF"),
            text=metas_periodo['atingimento'].apply(lambda x: f'{x:.1f}%'),
            textposition='top center',
            yaxis="y2",
            hovertemplate='<b>%{x}</b><br>Atingimento: %{y:.1f}%<extra></extra>'
        )
    )
    
    # Personalizar layout
    fig.update_layout(
        title={
            'text': "Atingimento de Metas por Período",
            'y':0.95,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': dict(family="Orbitron", size=24, color="#F8F8FF")
        },
        paper_bgcolor='rgba(13, 13, 13, 0.0)',
        plot_bgcolor='rgba(13, 13, 13, 0.0)',
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="right",
            x=1,
            font=dict(family="Montserrat", color="#F8F8FF")
        ),
        margin=dict(l=20, r=20, t=80, b=20),
        height=500,
        hovermode="x unified",
        yaxis=dict(
            title=dict(
                text="Faturamento (R$)",
                font=dict(family="Montserrat", color="#00FFFF")
            ),
            tickfont=dict(family="Montserrat", color="#00FFFF"),
            showgrid=True,
            gridcolor='rgba(0, 255, 255, 0.1)',
            zeroline=False
        ),
        yaxis2=dict(
            title=dict(
                text="Atingimento (%)",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            overlaying="y",
            side="right",
            showgrid=False,
            zeroline=False,
            range=[0, max(metas_periodo['atingimento']) * 1.2]
        ),
        xaxis=dict(
            title=dict(
                text="Período",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            showgrid=True,
            gridcolor='rgba(248, 248, 255, 0.1)',
            zeroline=False
        )
    )
    
    st.plotly_chart(fig, use_container_width=True)
    
    # Análise de tendência mensal
    st.markdown('<div class="section-title">Tendência Mensal por Categoria</div>', unsafe_allow_html=True)
    
    # Agrupar por categoria e período
    categoria_periodo = vendas.copy()
    if filtro_periodo:
        data_inicio, data_fim = filtro_periodo
        categoria_periodo = categoria_periodo[(categoria_periodo['data_venda'] >= data_inicio) & (categoria_periodo['data_venda'] <= data_fim)]
    
    if filtro_categorias:
        categoria_periodo = categoria_periodo[categoria_periodo['categoria'].isin(filtro_categorias)]
    
    if filtro_canais:
        categoria_periodo = categoria_periodo[categoria_periodo['canal_venda'].isin(filtro_canais)]
    
    categoria_periodo = categoria_periodo.groupby(['categoria', 'periodo']).agg({
        'id_venda': 'count'
    }).reset_index()
    
    # Criar gráfico de linha
    fig = go.Figure()
    
    # Cores para cada categoria
    cores = {
        'SUV': '#00FFFF',
        'Sedan': '#FF5F1F',
        'Hatch': '#F8F8FF',
        'Pickup': '#9933FF',
        'Elétrico': '#33FF99'
    }
    
    # Adicionar linha para cada categoria
    for categoria in sorted(categoria_periodo['categoria'].unique()):
        df_cat = categoria_periodo[categoria_periodo['categoria'] == categoria]
        
        fig.add_trace(
            go.Scatter(
                x=df_cat['periodo'],
                y=df_cat['id_venda'],
                name=categoria,
                line=dict(color=cores.get(categoria, '#F8F8FF'), width=3),
                mode='lines+markers',
                marker=dict(size=8, symbol='circle', color=cores.get(categoria, '#F8F8FF')),
                hovertemplate='<b>%{x}</b><br>Categoria: ' + categoria + '<br>Vendas: %{y:,}<extra></extra>'
            )
        )
    
    # Personalizar layout
    fig.update_layout(
        title={
            'text': "Evolução de Vendas por Categoria",
            'y':0.95,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': dict(family="Orbitron", size=24, color="#F8F8FF")
        },
        paper_bgcolor='rgba(13, 13, 13, 0.0)',
        plot_bgcolor='rgba(13, 13, 13, 0.0)',
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="right",
            x=1,
            font=dict(family="Montserrat", color="#F8F8FF")
        ),
        margin=dict(l=20, r=20, t=80, b=20),
        height=500,
        hovermode="x unified",
        yaxis=dict(
            title=dict(
                text="Quantidade de Vendas",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            showgrid=True,
            gridcolor='rgba(248, 248, 255, 0.1)',
            zeroline=False
        ),
        xaxis=dict(
            title=dict(
                text="Período",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            showgrid=True,
            gridcolor='rgba(248, 248, 255, 0.1)',
            zeroline=False
        )
    )
    
    st.plotly_chart(fig, use_container_width=True)
    
    # Download de dados
    st.markdown('<div class="download-section">', unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown(get_download_link(matriz_vendas.reset_index(), "matriz_vendas.csv", "📥 Baixar Dados de Vendas por Modelo"), unsafe_allow_html=True)
    
    with col2:
        st.markdown(get_download_link(metas_periodo, "metas_periodo.csv", "📥 Baixar Dados de Atingimento"), unsafe_allow_html=True)
    
    with col3:
        st.markdown(get_download_link(categoria_periodo, "categoria_periodo.csv", "📥 Baixar Dados de Tendência"), unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)

# Tab 5: Análise por Horário
with tab5:
    st.markdown('<div class="tab-title">Análise por Horário</div>', unsafe_allow_html=True)
    
    # Gráfico de dispersão por hora
    st.markdown('<div class="chart-container">', unsafe_allow_html=True)
    
    fig_scatter, vendas_hora = create_scatter_chart(vendas, filtro_periodo, filtro_categorias, filtro_canais)
    st.plotly_chart(fig_scatter, use_container_width=True)
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Análise por dia da semana
    st.markdown('<div class="section-title">Análise por Dia da Semana</div>', unsafe_allow_html=True)
    
    fig_line, vendas_dia = create_line_chart(vendas, filtro_periodo, filtro_categorias, filtro_canais)
    st.plotly_chart(fig_line, use_container_width=True)
    
    # Análise cruzada: Dia da Semana x Hora
    st.markdown('<div class="section-title">Análise Cruzada: Dia da Semana x Hora</div>', unsafe_allow_html=True)
    
    # Agrupar por dia da semana e hora
    dia_hora = vendas.copy()
    if filtro_periodo:
        data_inicio, data_fim = filtro_periodo
        dia_hora = dia_hora[(dia_hora['data_venda'] >= data_inicio) & (dia_hora['data_venda'] <= data_fim)]
    
    if filtro_categorias:
        dia_hora = dia_hora[dia_hora['categoria'].isin(filtro_categorias)]
    
    if filtro_canais:
        dia_hora = dia_hora[dia_hora['canal_venda'].isin(filtro_canais)]
    
    # Adicionar dia da semana
    dia_hora['dia_semana'] = dia_hora['data_venda'].dt.dayofweek
    dia_hora['dia_semana_nome'] = dia_hora['data_venda'].dt.day_name()
    
    # Mapeamento para nomes em português
    mapa_dias = {
        'Monday': 'Segunda-feira',
        'Tuesday': 'Terça-feira',
        'Wednesday': 'Quarta-feira',
        'Thursday': 'Quinta-feira',
        'Friday': 'Sexta-feira',
        'Saturday': 'Sábado',
        'Sunday': 'Domingo'
    }
    
    dia_hora['dia_semana_nome'] = dia_hora['dia_semana_nome'].map(mapa_dias)
    
    # Agrupar por dia da semana e hora
    dia_hora_stats = dia_hora.groupby(['dia_semana', 'dia_semana_nome', 'hora']).agg({
        'id_venda': 'count',
        'preco_venda': 'sum',
        'lucro': 'sum'
    }).reset_index()
    
    # Calcular margem
    dia_hora_stats['margem'] = (dia_hora_stats['lucro'] / dia_hora_stats['preco_venda']) * 100
    
    # Pivotar para criar matriz para heatmap
    matriz_dia_hora = dia_hora_stats.pivot_table(
        index='dia_semana_nome',
        columns='hora',
        values='id_venda',
        aggfunc='sum'
    ).fillna(0)
    
    # Ordenar por dia da semana
    ordem_dias = [mapa_dias[dia] for dia in ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']]
    matriz_dia_hora = matriz_dia_hora.reindex(ordem_dias)
    
    # Criar heatmap
    fig = go.Figure(data=go.Heatmap(
        z=matriz_dia_hora.values,
        x=matriz_dia_hora.columns,
        y=matriz_dia_hora.index,
        colorscale=[
            [0, 'rgba(13, 13, 13, 0.7)'],
            [0.2, 'rgba(255, 95, 31, 0.3)'],
            [0.4, 'rgba(255, 95, 31, 0.5)'],
            [0.6, 'rgba(255, 95, 31, 0.7)'],
            [0.8, 'rgba(255, 95, 31, 0.9)'],
            [1, '#FF5F1F']
        ],
        hovertemplate='<b>Dia:</b> %{y}<br><b>Hora:</b> %{x}:00<br><b>Vendas:</b> %{z}<extra></extra>',
        text=matriz_dia_hora.values.astype(int),
        texttemplate="%{text}",
        textfont={"family": "Montserrat", "size": 12, "color": "#F8F8FF"}
    ))
    
    # Personalizar layout
    fig.update_layout(
        title={
            'text': "Mapa de Calor: Vendas por Dia da Semana e Hora",
            'y':0.95,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': dict(family="Orbitron", size=24, color="#F8F8FF")
        },
        paper_bgcolor='rgba(13, 13, 13, 0.0)',
        plot_bgcolor='rgba(13, 13, 13, 0.0)',
        margin=dict(l=20, r=20, t=80, b=20),
        height=500,
        xaxis=dict(
            title=dict(
                text="Hora do Dia",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            showgrid=False,
            zeroline=False,
            tickmode='array',
            tickvals=list(range(0, 24)),
            ticktext=[f'{h}:00' for h in range(0, 24)]
        ),
        yaxis=dict(
            title=dict(
                text="Dia da Semana",
                font=dict(family="Montserrat", color="#F8F8FF")
            ),
            tickfont=dict(family="Montserrat", color="#F8F8FF"),
            showgrid=False,
            zeroline=False
        )
    )
    
    st.plotly_chart(fig, use_container_width=True)
    
    # Download de dados
    st.markdown('<div class="download-section">', unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown(get_download_link(vendas_hora, "vendas_hora.csv", "📥 Baixar Dados por Hora"), unsafe_allow_html=True)
    
    with col2:
        st.markdown(get_download_link(vendas_dia, "vendas_dia.csv", "📥 Baixar Dados por Dia da Semana"), unsafe_allow_html=True)
    
    with col3:
        st.markdown(get_download_link(dia_hora_stats, "dia_hora_stats.csv", "📥 Baixar Dados Cruzados"), unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)

# Tab 6: IA Insights
with tab6:
    st.markdown('<div class="tab-title">Análise Inteligente com IA</div>', unsafe_allow_html=True)
    
    # Insights gerados por IA
    st.markdown('<div class="section-title">Insights Gerados por IA</div>', unsafe_allow_html=True)
    
    narrativa = generate_narrative(insights_data)
    st.markdown(narrativa, unsafe_allow_html=True)
    
    # Recomendações estratégicas
    st.markdown('<div class="section-title">Recomendações Estratégicas</div>', unsafe_allow_html=True)
    
    recomendacoes = generate_strategic_recommendations(insights_data)
    st.markdown(recomendacoes, unsafe_allow_html=True)
    
    # Exportação de relatórios
    st.markdown('<div class="section-title">Exportação de Relatórios</div>', unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown('<div class="export-card">', unsafe_allow_html=True)
        st.markdown('<div class="export-title">Exportar para PDF</div>', unsafe_allow_html=True)
        st.markdown('<div class="export-description">Gere um relatório completo em PDF com todos os insights e recomendações.</div>', unsafe_allow_html=True)
        
        if st.button("Gerar PDF", key="pdf_button"):
            with st.spinner("Gerando relatório em PDF..."):
                pdf_buffer = export_to_pdf(vendas, metas, insights_data, filtro_periodo, filtro_categorias, filtro_canais)
                
                # Converter para base64 para download
                b64_pdf = base64.b64encode(pdf_buffer.read()).decode()
                href = f'<a href="data:application/pdf;base64,{b64_pdf}" download="relatorio_vendas.pdf" class="download-button">📥 Baixar Relatório PDF</a>'
                st.markdown(href, unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col2:
        st.markdown('<div class="export-card">', unsafe_allow_html=True)
        st.markdown('<div class="export-title">Exportar para PowerPoint</div>', unsafe_allow_html=True)
        st.markdown('<div class="export-description">Crie uma apresentação em PowerPoint com os principais insights e gráficos.</div>', unsafe_allow_html=True)
        
        if st.button("Gerar PowerPoint", key="ppt_button"):
            with st.spinner("Gerando apresentação em PowerPoint..."):
                ppt_buffer = export_to_ppt(vendas, metas, insights_data, filtro_periodo, filtro_categorias, filtro_canais)
                
                # Converter para base64 para download
                b64_ppt = base64.b64encode(ppt_buffer.read()).decode()
                href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="apresentacao_vendas.pptx" class="download-button">📥 Baixar Apresentação PowerPoint</a>'
                st.markdown(href, unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
